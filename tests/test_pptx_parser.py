"""Tests for PPTX input parser."""
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from backend.input_parser.pptx_parser import parse_pptx


def _create_sample_pptx(path: Path) -> None:
    """Create a sample PPTX with a title slide and a content slide."""
    prs = Presentation()

    # Slide 1: Title slide
    title_layout = prs.slide_layouts[0]  # Title Slide layout
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Q1 Review"
    slide1.placeholders[1].text = "Algorithm Team"

    # Slide 2: Content slide with bullet points
    content_layout = prs.slide_layouts[1]  # Title and Content layout
    slide2 = prs.slides.add_slide(content_layout)
    slide2.shapes.title.text = "Key Metrics"

    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "GPU utilization up 15%"

    p1 = tf.add_paragraph()
    p1.text = "Model accuracy improved to 97.3%"
    p1.level = 1

    p2 = tf.add_paragraph()
    p2.text = "Latency reduced by 20ms"
    p2.level = 1

    prs.save(str(path))


def test_parse_pptx_extracts_slides():
    """Check title='Q1 Review' and slide_count=2."""
    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = Path(tmp) / "test.pptx"
        _create_sample_pptx(pptx_path)

        result = parse_pptx(pptx_path)

        assert result["title"] == "Q1 Review"
        assert result["slide_count"] == 2
        assert len(result["slides"]) == 2
        assert result["slides"][0]["title"] == "Q1 Review"
        assert result["slides"][1]["title"] == "Key Metrics"
        assert result["slide_width"] is not None
        assert result["slide_height"] is not None


def test_parse_pptx_extracts_text():
    """Check second slide text_content contains 'GPU utilization'."""
    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = Path(tmp) / "test.pptx"
        _create_sample_pptx(pptx_path)

        result = parse_pptx(pptx_path)

        slide2 = result["slides"][1]
        assert "GPU utilization" in slide2["text_content"]
        assert "accuracy" in slide2["text_content"]
        assert "Latency" in slide2["text_content"]
        assert len(slide2["bullet_points"]) >= 2
