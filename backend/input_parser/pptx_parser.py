"""Parse PPTX files into structured data for PresentationIntent construction."""
from pathlib import Path
from pptx import Presentation as PptxPresentation


def parse_pptx(pptx_path: Path) -> dict:
    """Extract structured content from a PPTX file.

    Returns dict with title, slide_count, slide_width, slide_height, and per-slide data.
    """
    prs = PptxPresentation(str(pptx_path))
    title = ""
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "index": i,
            "text_content": "",
            "title": "",
            "bullet_points": [],
            "has_chart": False,
            "has_image": False,
            "has_table": False,
            "speaker_notes": "",
        }

        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
                        if paragraph.level > 0:
                            slide_data["bullet_points"].append(text)
            if shape.has_chart:
                slide_data["has_chart"] = True
            if shape.shape_type == 13:  # Picture
                slide_data["has_image"] = True
            if shape.has_table:
                slide_data["has_table"] = True

        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()
            if i == 0:
                title = slide_data["title"]

        slide_data["text_content"] = "\n".join(texts)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data["speaker_notes"] = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(slide_data)

    return {
        "title": title,
        "slide_count": len(slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": slides,
    }
